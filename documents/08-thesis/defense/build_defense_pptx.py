# -*- coding: utf-8 -*-
"""
build_defense_pptx.py — Dựng slide bảo vệ khóa luận (20 slide + phụ lục)
trên theme của template UTC `Presentation - File baocaodatn 16-9.pptx`.

Nguồn content: defense/defense-deck.html (40 slide) cô đọng 40->20 theo
defense/defense-deck-20slide-plan.md §2. Speaker notes (văn nói) nhúng vào
notes của mỗi slide.

Chạy:  .venv/bin/python defense/build_defense_pptx.py
Output: defense/KiteHub-baove-khoaluan-20slide.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
THESIS = os.path.dirname(HERE)
TEMPLATE = os.path.join(THESIS, "Presentation - File baocaodatn 16-9.pptx")
OUT = os.path.join(HERE, "KiteHub-baove-khoaluan-20slide.pptx")

# Theme palette (lấy từ template master)
NAVY = RGBColor(0x1F, 0x49, 0x7D)
BLUE = RGBColor(0x4F, 0x81, 0xBD)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xC0, 0x50, 0x4D)
GREY = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF3, 0xFA)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation(TEMPLATE)

# ---- xóa 3 slide mẫu trong template (giữ master/layout/theme) ----
# Dùng drop_rel để gỡ cả relationship + part (tránh trùng partname slideN.xml)
from pptx.oxml.ns import qn
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)

L_TITLE = prs.slide_layouts[0]      # Title Slide
L_CONTENT = prs.slide_layouts[1]    # Title and Content
L_SECTION = prs.slide_layouts[2]    # Section Header
L_TITLEONLY = prs.slide_layouts[5]  # Title Only
L_BLANK = prs.slide_layouts[6]      # Blank


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def style_title(slide, size=30):
    """Style placeholder title nếu có."""
    try:
        ph = slide.shapes.title
        if ph is not None:
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.bold = True
                    r.font.color.rgb = NAVY
    except Exception:
        pass


def add_title_only(title, size=28):
    s = prs.slides.add_slide(L_TITLEONLY)
    s.shapes.title.text = title
    style_title(s, size)
    return s


def add_bullets_box(slide, items, left=0.7, top=1.7, width=11.9, height=5.2,
                    base=18):
    """items: list of (text, level, bold, color|None)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        text, level, bold, color = (it + (None,))[:4] if len(it) < 4 else it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = ("• " if level == 0 else "– ") + text if level <= 1 else text
        run.font.size = Pt(base - level * 2)
        run.font.bold = bool(bold)
        run.font.color.rgb = color if color else RGBColor(0x22, 0x22, 0x22)
    return tb


def add_table(slide, rows, top=1.7, left=0.6, width=12.1, col_widths=None,
              fontsize=12, header_fontsize=12):
    nrows = len(rows)
    ncols = len(rows[0])
    height = Inches(0.4 * nrows)
    gtbl = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                                  Inches(width), height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gtbl.columns[i].width = Inches(w)
    for r in range(nrows):
        for c in range(ncols):
            cell = gtbl.cell(r, c)
            cell.text = rows[r][c]
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(header_fontsize if r == 0 else fontsize)
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
    return gtbl


def add_code_box(slide, code, top=1.7, left=0.7, width=11.9, height=5.0,
                 size=12):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x39)
    first = True
    for line in code.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xE6, 0xED, 0xF3)
    return tb


def add_caption(slide, text, top=6.9, size=11):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9),
                                  Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.italic = True
    r.font.color.rgb = GREY
    return tb


# ════════════════════════ S1 — Bìa ════════════════════════
s = prs.slides.add_slide(L_TITLE)
s.shapes.title.text = "Xây dựng hệ thống SaaS cung cấp dịch vụ đào tạo"
style_title(s, 32)
sub = s.placeholders[1]
sub.text = ("Nền tảng KiteHub & KiteClass — kiến trúc đa-tenant tích hợp AI\n"
            "\n"
            "Sinh viên: Nguyễn Văn Kiệt — MSSV 221230890 — Lớp CNTT1-K63\n"
            "Giảng viên hướng dẫn: TS. Nguyễn Đức Dư\n"
            "Bộ môn Công nghệ phần mềm — Khoa Công nghệ thông tin\n"
            "Trường Đại học Giao thông Vận tải — Hà Nội, 2026")
for p in sub.text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
set_notes(s, "Chào hội đồng. Em là Nguyễn Văn Kiệt, sinh viên lớp CNTT1-K63 "
             "dưới hướng dẫn của thầy Nguyễn Đức Dư. Hôm nay em xin trình bày "
             "khóa luận tốt nghiệp với đề tài Xây dựng hệ thống SaaS cung cấp "
             "dịch vụ đào tạo, gọi tắt là nền tảng KiteHub. (~30 giây)")

# ════════════════════════ S2 — Nội dung trình bày ════════════════════════
s = add_title_only("Nội dung trình bày")
add_bullets_box(s, [
    ("Phần 1 — Bối cảnh, vấn đề và mục tiêu đề tài", 0, True, NAVY),
    ("Phần 2 — Khảo sát thị trường và kỹ thuật AI tích hợp", 0, True, NAVY),
    ("Phần 3 — Phân tích, thiết kế kiến trúc đa-tenant", 0, True, NAVY),
    ("Phần 4 — Cài đặt, triển khai và kết quả đánh giá", 0, True, NAVY),
    ("Phần 5 — Demo trực tiếp, hạn chế và kết luận", 0, True, NAVY),
], top=2.0, base=22)
set_notes(s, "Bài trình bày gồm 5 phần: bối cảnh và mục tiêu, khảo sát thị "
             "trường và AI, thiết kế kiến trúc, triển khai và kết quả, cuối "
             "cùng là demo trực tiếp và kết luận. (~20 giây)")

# ════════════════════════ S3 — Bối cảnh & vấn đề ════════════════════════
s = add_title_only("Bối cảnh và vấn đề thực tế")
add_bullets_box(s, [
    ("Thông tư 29/2024/TT-BGDĐT chính thức hóa hoạt động dạy thêm có thu phí", 0, False, None),
    ("Ước tính hơn 50.000 trung tâm dạy thêm tư nhân đang hoạt động tại Việt Nam", 0, False, None),
    ("Phần lớn trung tâm vừa và nhỏ vẫn quản lý bằng Excel và nhóm Zalo", 0, False, None),
    ("Ba quan sát thúc đẩy đề tài:", 0, True, NAVY),
    ("Khoảng trống thị trường: trung tâm 100–2000 học viên cần phần mềm 500k–1,5tr đồng/tháng, tiếng Việt, đa-tenant gốc", 1, False, None),
    ("Mốc pháp lý: Luật Bảo vệ dữ liệu cá nhân 2023 hiệu lực 07/2026 — cơ hội tích hợp tuân thủ ngay từ thiết kế", 1, False, None),
    ("Công nghệ AI trưởng thành: API thương mại lần đầu cho phép tự động hóa nhận diện thương hiệu ở chi phí thấp", 1, False, None),
], top=1.6, base=17)
add_caption(s, "Nguồn: Magenest EdTech 2024; 6Wresearch Vietnam EdTech 2024–2030; VECITA 2024.")
set_notes(s, "Bối cảnh thúc đẩy đề tài: thị trường lớn, văn bản pháp luật hợp "
             "pháp hóa ngành dạy thêm, nhưng phần lớn trung tâm nhỏ vẫn quản lý "
             "thủ công vì phần mềm hiện có quá phức tạp hoặc quá đắt. Ba quan "
             "sát hội tụ: nhu cầu thị trường, mốc pháp lý, và công nghệ AI vừa "
             "đủ chín. (~60 giây)")

# ════════════════════════ S4 — Mục tiêu & phạm vi ════════════════════════
s = add_title_only("Mục tiêu và phạm vi đề tài")
add_bullets_box(s, [
    ("Bốn nhóm mục tiêu:", 0, True, NAVY),
    ("Chức năng: onboarding wizard, AI Branding, vòng đời tenant, tuân thủ tích hợp sẵn", 1, False, None),
    ("Phi chức năng: p95 ≤ 500ms (API đọc), cô lập cấp database, baseline OWASP Top 10, ≥100 tenant/instance", 1, False, None),
    ("Pháp lý: PDPL 2023 + Luật An ninh mạng 2018 + Thông tư 78/2021/TT-BTC", 1, False, None),
    ("Phương pháp luận: phát triển hướng chất lượng — kiểm thử trước, vòng lặp ngắn 1–3 ngày", 1, False, None),
    ("Phạm vi thực hiện: wizard tự phục vụ, AI Branding, cô lập RLS, tuân thủ pháp lý, triển khai AWS thực tế", 0, True, NAVY),
    ("Phát triển sau: thanh toán đa cổng, hóa đơn điện tử end-to-end, ứng dụng di động, mở rộng khối K-12", 0, True, GREY),
], top=1.6, base=16)
set_notes(s, "Bốn nhóm mục tiêu: không chỉ làm sản phẩm chạy được mà còn phải "
             "đo lường được, tuân thủ pháp luật, và có phương pháp luận chứng "
             "minh được. Phạm vi tập trung mức sẵn sàng cho tenant thực tế "
             "dùng; các tính năng nâng cao thuộc hướng phát triển sau. (~60 giây)")

# ════════════════════════ S5 — Khảo sát hệ thống tương tự ════════════════════════
s = add_title_only("Khảo sát hệ thống tương tự và yếu tố khác biệt")
add_table(s, [
    ["Sản phẩm", "Phân khúc", "Giá tham khảo", "Đa-tenant", "AI tích hợp"],
    ["BeeClass", "Trung tâm tiếng Anh", "1–3 tr/tháng", "Đơn tenant", "Không"],
    ["MISA AMIS Trường học", "Trường công K-12", "2–5 tr + setup", "Đơn tenant", "Không"],
    ["Mona eLMS", "Ngoại ngữ, tin học", "1–2 tr/tháng", "Đơn tenant", "Không"],
    ["Easy Edu", "Trung tâm K-12, kỹ năng", "0,8–2 tr/tháng", "Đơn tenant", "Không"],
    ["KiteHub (đề tài)", "Vừa & nhỏ 100–2000 HV", "0,5–1,5 tr/tháng", "Đa-tenant RLS gốc", "AI Branding"],
], top=1.7, fontsize=12, col_widths=[2.3, 2.6, 2.3, 2.4, 2.5])
add_caption(s, "Khoảng trống: trung tâm 100–2000 học viên cần đa-tenant gốc + AI tích hợp ở mức giá phải chăng.")
set_notes(s, "Khảo sát bốn sản phẩm tham khảo phổ biến tại Việt Nam. Hầu hết "
             "là đơn tenant, không có AI tích hợp. KiteHub khai thác khoảng "
             "trống đa-tenant gốc kết hợp AI Branding ở mức giá tier thấp. Chi "
             "tiết bốn phương pháp sinh ảnh để ở phụ lục. (~70 giây)")

# ════════════════════════ S6 — ⭐ AI Branding ════════════════════════
s = add_title_only("Đóng góp 1 — Kỹ thuật AI Branding tự động")
add_bullets_box(s, [
    ("Bài toán: sinh bộ nhận diện thương hiệu (logo, ảnh bìa, banner) tự động khi tenant đăng ký", 0, False, None),
    ("Quyết định kiến trúc: dùng API thương mại (Stable Diffusion XL qua Replicate) thay vì tự host GPU", 0, False, None),
    ("Lý do: chi phí GPU tự host 500–1000 USD/tháng; mô hình SOTA thay đổi mỗi 3–6 tháng", 1, False, None),
    ("Pipeline 6 bước (bất đồng bộ, không chặn giao diện):", 0, True, NAVY),
    ("Form chủ trung tâm → Gateway (xác thực) → Branding orchestrator → Hàng đợi RabbitMQ", 1, False, None),
    ("Worker xử lý async → AI provider (Replicate; dự phòng Hugging Face) → Quality Gate (lọc NSFW + brand-fit)", 1, False, None),
    ("Đạt → lưu MinIO/S3 theo tenant + thông báo chủ trung tâm; Không đạt → retry tối đa 3 lần", 1, False, None),
    ("Chi phí ~0,0036 USD cho 3 ảnh; thời gian sinh 30–60 giây", 0, True, GREEN),
], top=1.55, base=15)
set_notes(s, "Đóng góp thứ nhất: tự động hóa nhận diện thương hiệu. Quyết định "
             "không tự host AI là quyết định kiến trúc quan trọng, phù hợp quy "
             "mô và ngân sách. Pipeline gồm 6 bước, worker xử lý bất đồng bộ "
             "nên không chặn giao diện; có quality gate lọc nội dung và cơ chế "
             "dự phòng provider khi bị giới hạn. (~80 giây)")

# ════════════════════════ S7 — Tuân thủ pháp luật VN ════════════════════════
s = add_title_only("Tuân thủ pháp luật Việt Nam — 3 trụ cột")
add_table(s, [
    ["Văn bản", "Yêu cầu chính", "Cách hiện thực"],
    ["PDPL 2023\n(Luật 49/2023/QH15)", "Consent cụ thể; quyền truy cập/xóa dữ liệu; DPO khi >10k chủ thể; báo vi phạm 72h", "Bảng consent_record; quy trình DSAR; audit log bất biến (V60)"],
    ["Luật An ninh mạng 2018\n+ Nghị định 53/2022", "Lưu trữ dữ liệu tại VN khi vượt ngưỡng 1 triệu người dùng", "Lộ trình chuyển AWS Hà Nội Local Zone / Viettel / VNG khi đạt ngưỡng"],
    ["Thông tư 78/2021/TT-BTC", "Hóa đơn điện tử bắt buộc, kết nối cơ quan thuế", "Tích hợp đối tác MISA MeInvoice (không tự xây engine VAT)"],
], top=1.8, fontsize=12, col_widths=[3.2, 4.6, 4.3])
add_caption(s, "Cách tiếp cận: tuân thủ ngay từ thiết kế (compliance-by-design), không vá về sau.")
set_notes(s, "Ba trụ cột pháp luật được tích hợp ngay từ thiết kế. PDPL là mốc "
             "cứng tháng 7 năm 2026. Luật An ninh mạng có ngưỡng kích hoạt; "
             "hiện chưa chạm ngưỡng nhưng có lộ trình. Thông tư 78 về hóa đơn "
             "điện tử dùng đối tác chuyên trách. Mapping đầy đủ các điều luật "
             "PDPL ở phụ lục. (~60 giây)")

# ════════════════════════ S8 — Kiến trúc C4 L1 ════════════════════════
s = add_title_only("Kiến trúc tổng thể — C4 Level 1")
add_bullets_box(s, [
    ("Người dùng → Cloudflare (DNS, CDN, chống DDoS) → AWS ALB (kết thúc HTTPS, ap-southeast-1)", 0, False, None),
    ("KiteHub — mặt phẳng điều khiển (control-plane):", 0, True, NAVY),
    ("Gateway, Subscription, Branding, Email, Platform, Admin — quản lý vòng đời tenant", 1, False, None),
    ("KiteClass — mặt phẳng dữ liệu (data-plane):", 0, True, NAVY),
    ("Core (nghiệp vụ giáo dục) + Frontend (giao diện tenant)", 1, False, None),
    ("Hạ tầng dùng chung: PostgreSQL 16 (RLS theo tenant), S3 (tài nguyên), AWS SES (email)", 0, False, None),
    ("Hai mặt phẳng tách trách nhiệm rõ ràng nhưng chia sẻ một database cô lập bằng RLS", 0, True, GREEN),
], top=1.7, base=17)
set_notes(s, "Kiến trúc tổng thể chia hai mặt phẳng: KiteHub là control-plane "
             "quản lý vòng đời tenant, KiteClass là data-plane phục vụ nghiệp "
             "vụ giáo dục. Cả hai dùng chung một PostgreSQL với RLS cô lập theo "
             "tenant. Sơ đồ trực quan có trong bản deck đầy đủ. (~70 giây)")

# ════════════════════════ S9 — ⭐ Multi-tenant cô lập ════════════════════════
s = add_title_only("Đóng góp 2 — Chiến lược cô lập đa-tenant")
add_table(s, [
    ["Mô hình", "Chi phí", "Cô lập", "Quy mô", "Quyết định"],
    ["Instance mỗi tenant", "Rất cao", "Tuyệt đối", "≤10 tenant", "Loại — không mở rộng"],
    ["Database mỗi tenant", "Cao", "Mạnh", "10–100 tenant", "Loại — IAM phức tạp"],
    ["Schema mỗi tenant", "Trung bình", "Khá", "100–1000 tenant", "Cân nhắc sau"],
    ["Row-Level Security", "Thấp", "DB engine ép buộc", "≥1000 tenant", "Áp dụng"],
], top=1.9, fontsize=13, col_widths=[2.9, 1.9, 2.6, 2.3, 2.4])
add_caption(s, "Chọn RLS: native Postgres, chi phí thấp, kiểm chứng tại Salesforce, Shopify, HubSpot.")
set_notes(s, "So sánh bốn mô hình cô lập. Chọn Row-Level Security vì chi phí "
             "thấp nhất và được chính database engine ép buộc — không phụ thuộc "
             "lập trình viên nhớ thêm điều kiện lọc. Đây là đóng góp thứ hai "
             "của đề tài. (~70 giây)")

# ════════════════════════ S10 — RLS implementation ════════════════════════
s = add_title_only("Cài đặt PostgreSQL Row-Level Security")
add_code_box(s, """-- 1) Bật RLS cho mọi bảng thuộc phạm vi tenant
ALTER TABLE students ENABLE ROW LEVEL SECURITY;

-- 2) Policy: tenant chỉ thấy dữ liệu của chính mình
CREATE POLICY tenant_isolation ON students
  FOR ALL
  USING (tenant_id = current_setting('app.current_tenant_id')::uuid);

-- 3) Tầng ứng dụng: set biến phiên qua HikariCP connection init
cfg.setConnectionInitSql(
    "SET app.current_tenant_id = '" + tenantId + "'");""", top=1.7, size=14, height=4.2)
add_caption(s, "Postgres ép policy ở MỖI truy vấn — lập trình viên quên thêm điều kiện lọc cũng không rò dữ liệu.")
set_notes(s, "RLS cài đặt hai lớp: lớp database tạo policy, lớp ứng dụng set "
             "biến phiên qua HikariCP. Mỗi truy vấn Postgres tự động đánh giá "
             "policy nên ngay cả khi quên lọc cũng không rò dữ liệu chéo "
             "tenant. (~70 giây)")

# ════════════════════════ S11 — Defense-in-depth ════════════════════════
s = add_title_only("Bảo mật nhiều lớp — Defense-in-depth")
add_table(s, [
    ["Lớp", "Cơ chế", "Vai trò"],
    ["1. Xác thực JWT", "Spring Security filter", "Xác thực danh tính + claim tenant"],
    ["2. @PreAuthorize", "Kiểm tra quyền cấp method", "Phân quyền theo vai trò"],
    ["3. Tenant interceptor", "HikariCP init SQL", "Set biến phiên tenant hiện tại"],
    ["4. Bộ lọc repository", "JPA criteria", "Lớp dự phòng phòng khi RLS cấu hình sai"],
    ["5. PostgreSQL RLS", "Database engine", "Phòng tuyến cuối — DB ép buộc"],
], top=1.8, fontsize=13, col_widths=[2.8, 3.6, 5.7])
add_caption(s, "Một lớp lỗi không gây rò; phải đồng thời thủng cả 5 lớp mới rò dữ liệu chéo tenant.")
set_notes(s, "Nguyên tắc nhiều lớp độc lập. RLS là lớp cuối cùng: ngay cả khi "
             "lập trình viên quên kiểm tra ở tầng ứng dụng, Postgres vẫn ép "
             "buộc. Đây là khác biệt quan trọng so với cách chỉ lọc ở tầng ứng "
             "dụng. (~60 giây)")

# ════════════════════════ S12 — Tech stack + services ════════════════════════
s = add_title_only("Công nghệ và phân chia dịch vụ")
add_bullets_box(s, [
    ("KiteHub — 6 microservice (lifecycle khác nhau):", 0, True, NAVY),
    ("gateway, platform, subscription, branding, email, admin", 1, False, None),
    ("KiteClass — modular monolith:", 0, True, NAVY),
    ("kiteclass-core (domain giáo dục gắn kết chặt) + kiteclass-frontend (Next.js)", 1, False, None),
    ("Backend: Java 21 LTS + Spring Boot 3.5, Spring Security 6, JPA/Hibernate 6, MapStruct", 0, False, None),
    ("Dữ liệu & hàng đợi: PostgreSQL 16 + RLS, Flyway 10, Redis 7, RabbitMQ 3.13, MinIO/S3", 0, False, None),
    ("Frontend: TypeScript 5.7 + Next.js 15, React 19, Tailwind + Shadcn UI, TanStack Query", 0, False, None),
    ("DevOps: Docker, Terraform, GitHub Actions OIDC, AWS ECR/SES/ALB, CloudWatch + CloudTrail", 0, False, None),
], top=1.6, base=15)
set_notes(s, "Kiến trúc lai: KiteHub tách 6 microservice vì vòng đời khác nhau "
             "(branding bất đồng bộ, email hàng đợi, subscription giao dịch). "
             "KiteClass là modular monolith vì domain giáo dục gắn kết chặt. "
             "Stack chọn theo bản LTS để giảm rủi ro nâng cấp. (~50 giây)")

# ════════════════════════ S13 — ⭐ Implementation highlight ════════════════════════
s = add_title_only("Đóng góp 3 — Trích cài đặt: JWT Auth Filter")
add_code_box(s, """// kitehub-gateway: xác thực JWT một lần ở edge
protected void doFilterInternal(...) {
  String token = extractBearerToken(request);
  if (token != null && jwtUtil.validate(token)) {
    Claims claims = jwtUtil.parseClaims(token);
    String tenantId = claims.get("tenantId", String.class);
    String role     = claims.get("role", String.class);

    // Truyền tenant context xuống downstream service
    request.setAttribute("X-Tenant-Id", tenantId);

    var auth = new UsernamePasswordAuthenticationToken(
        claims.getSubject(), null,
        List.of(new SimpleGrantedAuthority("ROLE_" + role)));
    SecurityContextHolder.getContext().setAuthentication(auth);
  }
  filterChain.doFilter(request, response);
}""", top=1.7, size=12.5, height=4.6)
add_caption(s, "Xác thực JWT một lần ở gateway, truyền tenantId + role xuống dịch vụ — không lặp lại ở mỗi service.")
set_notes(s, "Đoạn code đại diện một trong các trích dẫn trong báo cáo. Mẫu "
             "thiết kế: xác thực JWT ở gateway, không lặp ở mỗi dịch vụ; tenant "
             "context truyền qua header để dịch vụ set biến phiên trước truy "
             "vấn. (~70 giây)")

# ════════════════════════ S14 — AWS Singapore ════════════════════════
s = add_title_only("Triển khai thực tế — AWS Singapore")
add_bullets_box(s, [
    ("Vùng ap-southeast-1: 2× EC2 t3.micro (backend + core/frontend) + RDS PostgreSQL db.t3.micro + S3", 0, False, None),
    ("Cloudflare (DNS, CDN) → AWS ALB (HTTPS) → EC2; SES gửi email; ECR lưu Docker image", 0, False, None),
    ("Lý do chọn AWS:", 0, True, NAVY),
    ("Đăng ký ổn định cho người dùng Việt Nam; hệ sinh thái trưởng thành (ECR, SES, ALB, Secrets Manager)", 1, False, None),
    ("Free Tier 12 tháng cho phép chi phí hạ tầng gần như 0 ở giai đoạn hiện tại", 1, False, None),
    ("Chưa chạm ngưỡng PDPL Điều 28 (10k chủ thể) hay Nghị định 53 (1 triệu người dùng)", 1, False, None),
    ("Phát triển sau: chuyển AWS Hà Nội Local Zone hoặc Viettel/VNG Cloud khi đạt ngưỡng lưu trữ trong nước", 0, True, GREY),
], top=1.6, base=16)
set_notes(s, "Triển khai thực tế trên AWS Singapore. Lý do chọn: đăng ký ổn "
             "định, hệ sinh thái trưởng thành, Free Tier 12 tháng cho chi phí "
             "gần như 0. Hiện chưa vượt ngưỡng pháp lý, có lộ trình chuyển vùng "
             "rõ ràng khi cần. (~60 giây)")

# ════════════════════════ S15 — CI/CD + Observability ════════════════════════
s = add_title_only("CI/CD và giám sát vận hành")
add_bullets_box(s, [
    ("CI/CD — GitHub Actions:", 0, True, NAVY),
    ("OIDC role assumption — không lưu AWS access key tĩnh trong secrets", 1, False, None),
    ("Docker tag bất biến theo SHA commit; Terraform plan-trước-apply + xác nhận thủ công", 1, False, None),
    ("Mọi terraform apply phải có audit artifact kiểm tra trạng thái trước khi chạy", 1, False, None),
    ("Giám sát 3 lớp:", 0, True, NAVY),
    ("CloudTrail — audit mọi lời gọi AWS API (bật trước khi tạo tài nguyên production)", 1, False, None),
    ("CloudWatch — log JSON có cấu trúc + cảnh báo (CPU >80%, ALB 5xx >1%)", 1, False, None),
    ("Prometheus + Grafana — metric ứng dụng (độ trễ outbox, request, bộ nhớ JVM)", 1, False, None),
], top=1.6, base=16)
set_notes(s, "CI/CD áp dụng chuẩn hiện đại: artifact bất biến, OIDC thay cho "
             "key tĩnh, xác nhận thủ công như một điểm dừng nhận thức. Giám sát "
             "ba lớp độc lập; CloudTrail bắt buộc bật trước khi tạo tài nguyên "
             "để có audit baseline. (~50 giây)")

# ════════════════════════ S16 — Quality-Driven ════════════════════════
s = add_title_only("Phương pháp luận phát triển hướng chất lượng")
add_table(s, [
    ["Trụ cột", "Cách áp dụng"],
    ["Vòng lặp ngắn", "Mỗi chu kỳ 1–3 ngày với phạm vi và tiêu chí nghiệm thu rõ ràng"],
    ["Kiểm thử trước", "Unit + integration test viết trước/song song; mục tiêu coverage ≥70%"],
    ["Đánh giá theo kỳ (audit)", "Mỗi miền có chu kỳ audit: bảo mật / hiệu năng / nghiệp vụ / giao diện / vận hành"],
    ["Chuẩn hóa quy tắc", "Mỗi sai sót phát hiện → quy tắc + cơ chế kiểm tra tự động + kiểm thử lại"],
], top=1.9, fontsize=14, col_widths=[3.0, 9.2])
add_caption(s, "Cơ sở lý thuyết: Deming PDCA, Beck TDD, Poppendieck Lean, IEEE 730-2014 SQA.")
set_notes(s, "Bốn trụ cột phương pháp luận, có cơ sở lý thuyết Deming, Beck, "
             "Poppendieck, IEEE 730. Đặc biệt trụ cột thứ tư: mỗi sai sót trở "
             "thành một quy tắc có cơ chế kiểm tra, không chỉ ghi chú lần sau "
             "cẩn thận hơn. (~50 giây)")

# ════════════════════════ S17 — Kết quả & KPI ════════════════════════
s = add_title_only("Kết quả đánh giá — các chỉ số chính")
add_table(s, [
    ["Hạng mục", "Baseline", "Hiện tại", "Tăng", "Động lực cải tiến"],
    ["Hiệu năng", "81/100 B", "86/100 B+", "+5", "RLS NULL force-fail, reset GUC, phân trang con trỏ"],
    ["Bảo mật", "76/100 C+", "93/100 A", "+17", "Audit OWASP, audit_logs bất biến (V60), chặn bypass admin"],
    ["Chất lượng", "78/110 C+", "90/110 B+", "+12", "Coverage test, giám sát, sẵn sàng vận hành"],
    ["Tải (RPS)", "—", "≥100k", "—", "Kiểm thử tải Locust — dư địa cho cohort mời thử"],
], top=1.85, fontsize=12.5, col_widths=[2.0, 1.7, 1.7, 1.1, 5.7])
add_caption(s, "Mỗi điểm số gắn audit report có evidence block trong báo cáo Chương 4 — không phải tự nhận định.")
set_notes(s, "Bốn chỉ số chính: hiệu năng 86, bảo mật 93, chất lượng 90 trên "
             "110, đều vượt ngưỡng đạt. Quan trọng: mỗi điểm số có audit report "
             "với evidence block làm chứng cứ, và trajectory cho thấy cải tiến "
             "liên tục có thể đo được. (~80 giây)")

# ════════════════════════ S18 — DEMO ════════════════════════
s = prs.slides.add_slide(L_SECTION)
s.shapes.title.text = "Demo trực tiếp"
style_title(s, 36)
try:
    body = s.placeholders[1]
    body.text = ("Luồng: khách tham quan → đăng ký onboarding → wizard tạo "
                 "tenant → chứng minh cô lập đa-tenant\n"
                 "Dự phòng: video backup nếu có sự cố mạng/hạ tầng")
    for p in body.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(16)
except Exception:
    pass
set_notes(s, "Chuyển sang demo trực tiếp theo kịch bản 6 phase: khách tham "
             "quan trang công khai, đăng ký onboarding, wizard tạo tenant, "
             "chứng minh cô lập đa-tenant bằng hai tài khoản khác tenant, xem "
             "audit log. Nếu sự cố sẽ dùng video dự phòng. (~30 giây + demo)")

# ════════════════════════ S19 — Hạn chế & hướng phát triển ════════════════════════
s = add_title_only("Hạn chế thừa nhận và hướng phát triển")
add_bullets_box(s, [
    ("Hạn chế hiện tại:", 0, True, RED),
    ("Hạ tầng 2× t3.micro bộ nhớ hạn chế; chưa có failover đa vùng sẵn sàng (multi-AZ)", 1, False, None),
    ("Chưa lưu trữ dữ liệu trong nước — kích hoạt chuyển vùng khi vượt ngưỡng người dùng", 1, False, None),
    ("Quality gate AI dùng bộ lọc của provider; chưa tự huấn luyện bộ phân loại brand-fit", 1, False, None),
    ("Chưa có ứng dụng di động native — người dùng truy cập qua web responsive", 1, False, None),
    ("Hướng phát triển sau:", 0, True, GREEN),
    ("Mở đăng ký công khai, nâng cấp hạ tầng, kích hoạt cổng thanh toán MoMo/VNPay", 1, False, None),
    ("Tích hợp hóa đơn điện tử (MISA MeInvoice), kênh Zalo OA, mở rộng khối K-12 sau khi có legal counsel", 1, False, None),
], top=1.6, base=16)
set_notes(s, "Thừa nhận hạn chế trung thực kèm lộ trình là cách tiếp cận tốt "
             "hơn che giấu rồi bị hội đồng phát hiện. Mỗi hạn chế đều có hướng "
             "phát triển tương ứng: nâng cấp hạ tầng, chuyển vùng dữ liệu, mở "
             "thanh toán, tích hợp Zalo và hóa đơn điện tử qua đối tác. (~70 giây)")

# ════════════════════════ S20 — Kết luận ════════════════════════
s = add_title_only("Kết luận")
add_bullets_box(s, [
    ("Ba đóng góp chính của đề tài:", 0, True, NAVY),
    ("Kỹ thuật AI Branding tự động — sinh bộ nhận diện thương hiệu chi phí thấp, bất đồng bộ", 1, False, None),
    ("Kiến trúc đa-tenant cô lập bằng PostgreSQL RLS — bảo mật nhiều lớp, mở rộng tốt", 1, False, None),
    ("Phương pháp luận phát triển hướng chất lượng — audit có evidence, chuẩn hóa quy tắc", 1, False, None),
    ("Sản phẩm đã triển khai thực tế trên AWS, đạt các ngưỡng đánh giá hiệu năng/bảo mật/chất lượng", 0, False, None),
    ("Em xin chân thành cảm ơn thầy Nguyễn Đức Dư và quý hội đồng — sẵn sàng trả lời câu hỏi", 0, True, NAVY),
], top=1.8, base=18)
set_notes(s, "Tóm lại ba đóng góp: kỹ thuật AI Branding, kiến trúc đa-tenant "
             "RLS, và phương pháp luận hướng chất lượng. Sản phẩm đã triển khai "
             "thực tế và đạt các ngưỡng đánh giá. Cảm ơn GVHD và hội đồng, em "
             "sẵn sàng nhận câu hỏi. (~40 giây)")

# ════════════════════════ PHỤ LỤC ════════════════════════
s = prs.slides.add_slide(L_SECTION)
s.shapes.title.text = "Phụ lục — Slide dự phòng cho phần hỏi đáp"
style_title(s, 28)
set_notes(s, "Các slide phụ lục bật khi hội đồng hỏi sâu, theo 4 nhóm: Kiến "
             "trúc / Phi chức năng / Nghiệp vụ-Pháp lý / Quy trình.")

# A1 — 4 phương pháp text-to-image
s = add_title_only("Phụ lục A1 — So sánh 4 phương pháp sinh ảnh")
add_table(s, [
    ["Mô hình", "Provider", "Chi phí/ảnh", "Chất lượng", "Độ trễ", "Quyết định"],
    ["Stable Diffusion XL", "Replicate", "~0,0012 USD", "Cao", "4–8s", "Áp dụng chính"],
    ["SDXL Turbo", "Hugging Face", "Free tier", "Khá", "1–2s", "Dự phòng"],
    ["DALL-E 3", "OpenAI", "~0,04 USD", "Rất cao", "8–12s", "Đắt gấp ~30 lần"],
    ["Midjourney v6", "—", "—", "Rất cao", "—", "Không có API công khai"],
], top=1.9, fontsize=12.5, col_widths=[2.7, 2.2, 1.9, 1.8, 1.4, 2.2])
set_notes(s, "Chọn Stable Diffusion XL vì cân bằng chi phí, chất lượng, độ "
             "trễ. DALL-E 3 đắt gấp khoảng 30 lần, không phù hợp tier thấp.")

# A2 — PDPL mapping đầy đủ
s = add_title_only("Phụ lục A2 — PDPL 2023: điều luật → tính năng")
add_table(s, [
    ["Điều", "Yêu cầu", "Cách hiện thực"],
    ["Điều 9", "Quyền truy cập dữ liệu của chủ thể", "Endpoint DSAR + dashboard phụ huynh"],
    ["Điều 11", "Audit log bất biến, chống sửa", "Migration V60 — admin_audit_logs không cho UPDATE/DELETE"],
    ["Điều 16", "Consent cụ thể theo từng mục đích", "Bảng consent_record + luồng thu hồi"],
    ["Điều 28", "DPO khi >10k chủ thể", "Hiện <1000 chủ thể — chưa kích hoạt; có lộ trình"],
    ["Điều 30", "Báo vi phạm trong 72h", "Quy trình incident response + cảnh báo CloudWatch SNS"],
], top=1.9, fontsize=12.5, col_widths=[1.4, 4.6, 6.2])
set_notes(s, "Năm điều luật cốt lõi mapping 1-1 sang tính năng. Điều 11 audit "
             "log bất biến là yêu cầu khó nhất, giải bằng bảng không cho sửa/"
             "xóa ở cấp database.")

# A3 — Auth sequence (text)
s = add_title_only("Phụ lục A3 — Luồng xác thực và truy vấn")
add_bullets_box(s, [
    ("Người dùng đăng nhập (email + mật khẩu) → Frontend gọi POST /api/auth/login", 0, False, None),
    ("Gateway xác thực, dịch vụ trả JWT chứa tenantId + role; set HttpOnly cookie", 0, False, None),
    ("Request có cookie → Gateway verify JWT, trích tenantId, forward kèm header X-Tenant-Id", 0, False, None),
    ("Dịch vụ set app.current_tenant_id rồi truy vấn → Postgres tự lọc theo RLS policy", 0, False, None),
    ("Kết quả chỉ chứa dữ liệu của đúng tenant — không có WHERE tenant_id viết tay", 0, True, GREEN),
], top=1.8, base=17)
set_notes(s, "Toàn bộ luồng không có điều kiện lọc tenant viết tay ở repository "
             "— database tự ép buộc.")

# A4 — Cost breakdown
s = add_title_only("Phụ lục A4 — Phân tích chi phí")
add_table(s, [
    ["Hạng mục", "Chi phí hiện tại", "Ghi chú"],
    ["AWS Free Tier (12 tháng)", "0 USD/tháng", "2× t3.micro + RDS db.t3.micro + 5GB S3"],
    ["AWS SES email", "0 USD/tháng", "62k email/tháng miễn phí khi gửi từ EC2"],
    ["AI Branding (Replicate)", "~0,0036 USD/tenant", "3 ảnh × 0,0012 USD"],
    ["Cloudflare DNS + CDN", "0 USD/tháng", "Gói Free đủ dùng"],
    ["Tên miền kitehub.me", "~9 USD/năm", "Gia hạn hàng năm"],
], top=1.9, fontsize=13, col_widths=[3.4, 3.0, 5.8])
set_notes(s, "Chi phí hạ tầng gần như 0 nhờ Free Tier; chi phí chính là AI "
             "Branding khoảng 0,19 USD mỗi lần onboard tenant. Số liệu chính "
             "xác tổng hợp khi có dữ liệu vận hành thực tế.")

# A5 — Market segment / persona
s = add_title_only("Phụ lục A5 — Phân khúc và persona mục tiêu")
add_bullets_box(s, [
    ("Định vị: trung tâm dạy thêm vừa-nhỏ (1–10 chi nhánh, 100–2000 học viên)", 0, False, None),
    ("Persona chính hiện tại:", 0, True, NAVY),
    ("P1 — Giáo viên độc lập (1–50 học viên, tự quản lý)", 1, False, None),
    ("P2 — Chủ trung tâm (1–10 chi nhánh, 100–2000 học viên)", 1, False, None),
    ("P3 — Quản lý vận hành cấp chi nhánh", 1, False, None),
    ("Persona phát triển sau: phụ huynh, học viên (giao diện riêng)", 0, True, GREY),
], top=1.8, base=17)
set_notes(s, "Ba persona chính P1-P2-P3 là trọng tâm hiện tại; phụ huynh và "
             "học viên hỗ trợ qua giao diện riêng, thuộc hướng phát triển sau.")

# A6 — Cohort plan
s = add_title_only("Phụ lục A6 — Lộ trình mời tenant thực tế")
add_bullets_box(s, [
    ("Tuần 1–2: tiếp cận 10–15 trung tâm tiềm năng qua mạng lưới giáo viên", 0, False, None),
    ("Tuần 3–4: onboarding 4–5 trung tâm — chủ trung tâm tự setup qua wizard", 0, False, None),
    ("Tuần 5–7: vận hành thực tế — quản lý học viên, điểm danh, thanh toán", 0, False, None),
    ("Tuần 8–9: phỏng vấn feedback + thu xác nhận đánh giá có chữ ký, đóng cohort", 0, False, None),
    ("Mục tiêu: ≥4 trung tâm ký xác nhận đã dùng thực tế — phụ lục đánh giá Chương 4", 0, True, GREEN),
], top=1.8, base=17)
set_notes(s, "Lộ trình mời tenant thực tế nhằm đạt mục tiêu: ít nhất 4 trung "
             "tâm ký xác nhận đã sử dụng thực tế — khác biệt giữa thesis demo "
             "trên máy và thesis có người dùng thật xác nhận.")

prs.save(OUT)
print("OK saved:", OUT)
print("Tổng số slide:", len(prs.slides._sldIdLst))
